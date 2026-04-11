"""
文档加载器 —— RAG 文档摄取的第一层

这一层不要只理解成“读文件”。在更工程化的 RAG 项目里，loader 负责把
PDF/Word/Excel/PPT 等不同格式统一成结构化元素，后续 splitter、embedding、
retriever 才能拿到页码、表格、标题等元数据。
"""

from __future__ import annotations

import hashlib
import logging
import os
import re
from dataclasses import dataclass, field
from io import StringIO
from pathlib import Path
from typing import Any


logger = logging.getLogger(__name__)


@dataclass
class DocumentElement:
    """
    文档中的一个可检索元素。

    先把文档拆成“元素”，再在 splitter 中组合成 chunk，这是比直接拼接整篇
    文本更适合 RAG 的做法。这里保留 page/bbox/type，是为了后面能做页码溯源、
    结构感知分块和表格专项处理。
    """

    text: str
    element_type: str
    source: str
    page: int | None = None
    bbox: tuple[float, float, float, float] | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class LoadedDocument:
    """一次文档加载的结构化结果。"""

    doc_id: str
    source: str
    file_hash: str
    parser: str
    elements: list[DocumentElement]
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def text(self) -> str:
        """兼容旧流程：把结构化元素拼回纯文本。"""
        return "\n\n".join(element.text for element in self.elements if element.text.strip())


def load_document(filepath: str | os.PathLike[str]) -> LoadedDocument:
    """
    加载文档并返回结构化结果。

    当前阶段先完成 1-3：
    1. 建立 LoadedDocument / DocumentElement 结构
    2. PDF 默认使用 PyMuPDF，保留页码和 block 信息
    3. PyMuPDF 不可用或抽取失败时回退到 pdfminer.six
    """
    path = Path(filepath)
    file_ext = path.suffix.lower()
    file_hash = _calculate_file_hash(path)
    doc_id = _build_doc_id(file_hash)

    if file_ext == ".pdf":
        return _load_pdf(path, doc_id, file_hash)
    if file_ext in [".txt", ".md"]:
        return _load_plain_text(path, doc_id, file_hash)
    if file_ext == ".docx":
        return _load_docx(path, doc_id, file_hash)
    if file_ext in [".xlsx", ".xls"]:
        return _load_excel(path, doc_id, file_hash)
    if file_ext == ".pptx":
        return _load_pptx(path, doc_id, file_hash)

    logger.warning("不支持的文件格式: %s", file_ext)
    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="unsupported",
        elements=[],
        metadata={"file_ext": file_ext, "status": "unsupported"},
    )


def extract_text(filepath: str | os.PathLike[str]) -> str:
    """
    兼容旧接口：从文件中提取纯文本内容。

    现有 rag_demo.py 还在调用 extract_text()。为了避免一次重构影响整条链路，
    这里内部改为调用新的 load_document()，再把元素拼回字符串。
    """
    return load_document(filepath).text


def _load_pdf(path: Path, doc_id: str, file_hash: str) -> LoadedDocument:
    """PDF 主解析器：优先 PyMuPDF，失败或空文本时回退 pdfminer。"""
    try:
        loaded = _load_pdf_with_pymupdf(path, doc_id, file_hash)
        if loaded.elements:
            return loaded
        logger.warning("PyMuPDF 未抽取到文本，尝试 pdfminer fallback: %s", path.name)
    except ImportError:
        logger.warning("处理 PDF 建议安装 PyMuPDF，正在尝试 pdfminer fallback")
    except Exception as exc:
        logger.warning("PyMuPDF 解析失败，尝试 pdfminer fallback: %s", exc)

    return _load_pdf_with_pdfminer(path, doc_id, file_hash)


def _load_pdf_with_pymupdf(path: Path, doc_id: str, file_hash: str) -> LoadedDocument:
    """
    使用 PyMuPDF 解析 PDF。

    get_text("blocks") 会返回带坐标的文本块，比整页纯文本更适合后续做
    “第几页第几个区域”的来源追踪。
    """
    import fitz

    elements: list[DocumentElement] = []
    with fitz.open(path) as pdf:
        for page_index, page in enumerate(pdf, start=1):
            blocks = page.get_text("blocks")
            for block_no, block in enumerate(blocks):
                x0, y0, x1, y1, text, *_ = block
                cleaned = _clean_text(text)
                if not cleaned:
                    continue
                elements.append(
                    DocumentElement(
                        text=cleaned,
                        element_type="pdf_block",
                        source=path.name,
                        page=page_index,
                        bbox=(float(x0), float(y0), float(x1), float(y1)),
                        metadata={"block_no": block_no},
                    )
                )

        page_count = pdf.page_count

    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="pymupdf",
        elements=elements,
        metadata={"file_ext": ".pdf", "page_count": page_count, "status": "ok"},
    )


def _load_pdf_with_pdfminer(path: Path, doc_id: str, file_hash: str) -> LoadedDocument:
    """PDF fallback：保留旧项目的 pdfminer.six 能力，保证兼容性。"""
    try:
        from pdfminer.high_level import extract_text_to_fp
    except ImportError:
        logger.error("处理 PDF 需要安装 PyMuPDF 或 pdfminer.six")
        return _failed_document(path, doc_id, file_hash, ".pdf", "pdf_parser_missing")

    output = StringIO()
    try:
        with open(path, "rb") as file:
            extract_text_to_fp(file, output)
    except Exception as exc:
        logger.error("pdfminer 解析 PDF 失败: %s", exc)
        return _failed_document(path, doc_id, file_hash, ".pdf", "pdf_parse_failed")

    text = _clean_text(output.getvalue())
    elements = []
    if text:
        elements.append(
            DocumentElement(
                text=text,
                element_type="pdf_text",
                source=path.name,
                metadata={"fallback": True},
            )
        )

    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="pdfminer",
        elements=elements,
        metadata={"file_ext": ".pdf", "status": "ok" if elements else "empty"},
    )


def _load_plain_text(path: Path, doc_id: str, file_hash: str) -> LoadedDocument:
    with open(path, "r", encoding="utf-8") as file:
        text = _clean_text(file.read())

    elements = [
        DocumentElement(text=text, element_type="markdown" if path.suffix.lower() == ".md" else "text", source=path.name)
    ] if text else []

    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="plain_text",
        elements=elements,
        metadata={"file_ext": path.suffix.lower(), "status": "ok" if elements else "empty"},
    )


def _load_docx(path: Path, doc_id: str, file_hash: str) -> LoadedDocument:
    try:
        from docx import Document
    except ImportError:
        logger.error("处理 Word 文档需要安装 python-docx")
        return _failed_document(path, doc_id, file_hash, ".docx", "docx_parser_missing")

    doc = Document(path)
    elements: list[DocumentElement] = []
    for para_no, para in enumerate(doc.paragraphs):
        text = _clean_text(para.text)
        if not text:
            continue
        elements.append(
            DocumentElement(
                text=text,
                element_type="docx_paragraph",
                source=path.name,
                metadata={"paragraph_no": para_no, "style": para.style.name if para.style else None},
            )
        )

    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="python-docx",
        elements=elements,
        metadata={"file_ext": ".docx", "status": "ok" if elements else "empty"},
    )


def _load_excel(path: Path, doc_id: str, file_hash: str) -> LoadedDocument:
    try:
        import pandas as pd
    except ImportError:
        logger.error("处理 Excel 文件需要安装 pandas")
        return _failed_document(path, doc_id, file_hash, path.suffix.lower(), "excel_parser_missing")

    elements: list[DocumentElement] = []
    xl = pd.ExcelFile(path)
    for sheet_no, sheet_name in enumerate(xl.sheet_names):
        df = xl.parse(sheet_name)
        text = _clean_text(f"工作表: {sheet_name}\n{df.to_string(index=False)}")
        if not text:
            continue
        elements.append(
            DocumentElement(
                text=text,
                element_type="excel_sheet",
                source=path.name,
                metadata={"sheet_no": sheet_no, "sheet_name": sheet_name, "rows": len(df), "columns": len(df.columns)},
            )
        )

    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="pandas",
        elements=elements,
        metadata={"file_ext": path.suffix.lower(), "sheet_count": len(xl.sheet_names), "status": "ok" if elements else "empty"},
    )


def _load_pptx(path: Path, doc_id: str, file_hash: str) -> LoadedDocument:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("处理 PPT 文件需要安装 python-pptx")
        return _failed_document(path, doc_id, file_hash, ".pptx", "pptx_parser_missing")

    prs = Presentation(path)
    elements: list[DocumentElement] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape_no, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text"):
                continue
            text = _clean_text(shape.text)
            if not text:
                continue
            elements.append(
                DocumentElement(
                    text=text,
                    element_type="pptx_shape",
                    source=path.name,
                    page=slide_no,
                    metadata={"slide_no": slide_no, "shape_no": shape_no},
                )
            )

    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="python-pptx",
        elements=elements,
        metadata={"file_ext": ".pptx", "slide_count": len(prs.slides), "status": "ok" if elements else "empty"},
    )


def _calculate_file_hash(path: Path) -> str:
    """计算文件 hash：后续可用于去重、缓存和增量索引。"""
    sha256 = hashlib.sha256()
    with open(path, "rb") as file:
        for chunk in iter(lambda: file.read(1024 * 1024), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


def _build_doc_id(file_hash: str) -> str:
    return f"doc_{file_hash[:16]}"


def _clean_text(text: str) -> str:
    """轻量清洗：先只做空白归一化，避免过早引入复杂规则误删内容。"""
    if not text:
        return ""
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _failed_document(path: Path, doc_id: str, file_hash: str, file_ext: str, status: str) -> LoadedDocument:
    return LoadedDocument(
        doc_id=doc_id,
        source=path.name,
        file_hash=file_hash,
        parser="failed",
        elements=[],
        metadata={"file_ext": file_ext, "status": status},
    )
